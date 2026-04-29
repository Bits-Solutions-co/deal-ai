"""
Convert DEAL_AI_FUNCTIONAL_DOCUMENTATION.md to DOCX and PPTX.

Custom MD parser + python-docx / python-pptx renderers.
Handles: headings (#..####), paragraphs, bullet lists, code blocks, fenced code,
tables (pipe-separated), inline emphasis (**bold**, `code`), block quotes.

Mermaid diagrams render as preformatted code blocks in DOCX
and as a single labeled bullet on a slide in PPTX.
"""

import base64
import hashlib
import re
import urllib.request
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu as PEmu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE


HERE = Path(__file__).resolve().parent
MD_PATH = HERE / "DEAL_AI_FUNCTIONAL_DOCUMENTATION.md"
DOCX_PATH = HERE / "DEAL_AI_FUNCTIONAL_DOCUMENTATION.docx"
PPTX_PATH = HERE / "DEAL_AI_FUNCTIONAL_DOCUMENTATION.pptx"
MERMAID_CACHE = HERE / ".mermaid-cache"


def _sanitize_mermaid(source: str) -> str:
    """Replace Unicode glyphs and stateDiagram <br/> that mermaid.ink rejects."""
    s = source
    # Unicode arrows / comparisons → ASCII
    replacements = {
        "→": "->",
        "←": "<-",
        "↔": "<->",
        "▸": ">",
        "≤": "<=",
        "≥": ">=",
        "—": "-",
        "–": "-",
        "❝": '"',
        "❞": '"',
    }
    for k, v in replacements.items():
        s = s.replace(k, v)
    # stateDiagram-v2 transition labels don't support <br/> — convert to single line
    first = s.lstrip().split("\n", 1)[0]
    if first.startswith("stateDiagram"):
        s = re.sub(r"<br\s*/?>", " | ", s)
    return s


def _encode_base64(source: str) -> str:
    return base64.urlsafe_b64encode(source.encode("utf-8")).decode("ascii")


def _encode_pako(source: str) -> str:
    """zlib (deflate) + base64url, matching mermaid.ink's pako: scheme."""
    import zlib
    compressed = zlib.compress(source.encode("utf-8"), 9)
    return base64.urlsafe_b64encode(compressed).decode("ascii").rstrip("=")


def _try_fetch(url: str) -> bytes | None:
    req = urllib.request.Request(url, headers={"User-Agent": "convert-docs.py"})
    with urllib.request.urlopen(req, timeout=60) as resp:
        data = resp.read()
    if not data or data[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    return data


def fetch_mermaid_png(source: str) -> Path | None:
    """Render via mermaid.ink, trying pako-compressed first then plain base64.
    Caches PNGs by hash of source. Retries on transient errors."""
    import time
    MERMAID_CACHE.mkdir(exist_ok=True)
    h = hashlib.sha256(source.encode("utf-8")).hexdigest()[:16]
    cache_path = MERMAID_CACHE / f"{h}.png"
    if cache_path.exists() and cache_path.stat().st_size > 0:
        return cache_path

    sanitized = _sanitize_mermaid(source)

    # Mermaid.ink accepts both pako: (compressed) and raw base64 forms.
    # pako is shorter — preferred for big diagrams that exceed URL limits.
    schemes = [
        ("pako", f"https://mermaid.ink/img/pako:{_encode_pako(sanitized)}?type=png&bgColor=FFFFFF"),
        ("base64", f"https://mermaid.ink/img/{_encode_base64(sanitized)}?type=png&bgColor=FFFFFF"),
    ]

    for scheme_name, url in schemes:
        max_attempts = 3
        last_err = None
        for attempt in range(1, max_attempts + 1):
            try:
                data = _try_fetch(url)
                if data is None:
                    last_err = "non-PNG response"
                    break
                cache_path.write_bytes(data)
                return cache_path
            except urllib.error.HTTPError as e:
                last_err = f"HTTP {e.code} {e.reason}"
                if e.code in (503, 502, 504, 429) and attempt < max_attempts:
                    wait = 2 * attempt
                    time.sleep(wait)
                    continue
                break  # 400/404 etc — don't retry, try next scheme
            except Exception as e:
                last_err = str(e)
                if attempt < max_attempts:
                    time.sleep(2 * attempt)
                    continue
                break
        print(f"  mermaid {scheme_name} failed: {last_err}")

    preview_lines = source.strip().split("\n")[:6]
    print(f"    [unrenderable diagram, source first 6 lines]:")
    for ln in preview_lines:
        print(f"      | {ln[:160]}")
    return None


# ---------- Block parser -------------------------------------------------------

class Block:
    def __init__(self, kind, **kwargs):
        self.kind = kind
        self.__dict__.update(kwargs)


def parse_md(text: str):
    lines = text.splitlines()
    blocks = []
    i = 0
    while i < len(lines):
        line = lines[i]

        # Fenced code (``` or ```mermaid)
        if line.startswith("```"):
            lang = line[3:].strip()
            buf = []
            i += 1
            while i < len(lines) and not lines[i].startswith("```"):
                buf.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            blocks.append(Block("code", lang=lang, content="\n".join(buf)))
            continue

        # Heading
        m = re.match(r"^(#{1,6})\s+(.*)", line)
        if m:
            level = len(m.group(1))
            blocks.append(Block("heading", level=level, text=m.group(2).strip()))
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^---+\s*$", line):
            blocks.append(Block("hr"))
            i += 1
            continue

        # Table (header row followed by separator)
        if "|" in line and i + 1 < len(lines) and re.match(
            r"^\s*\|?\s*[:\- |]+\|[:\- |]+\s*$", lines[i + 1]
        ):
            header_cells = [c.strip() for c in line.strip().strip("|").split("|")]
            i += 2
            rows = []
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                rows.append(cells)
                i += 1
            blocks.append(Block("table", header=header_cells, rows=rows))
            continue

        # Block quote
        if line.startswith(">"):
            buf = []
            while i < len(lines) and lines[i].startswith(">"):
                buf.append(lines[i].lstrip(">").lstrip())
                i += 1
            blocks.append(Block("quote", content="\n".join(buf)))
            continue

        # Bullet list
        if re.match(r"^\s*[-*+]\s+", line) or re.match(r"^\s*\d+\.\s+", line):
            buf = []
            while i < len(lines) and (
                re.match(r"^\s*[-*+]\s+", lines[i])
                or re.match(r"^\s*\d+\.\s+", lines[i])
                or (lines[i].startswith("  ") and lines[i].strip())
            ):
                buf.append(lines[i])
                i += 1
            blocks.append(Block("list", items=buf))
            continue

        # Blank line — skip
        if not line.strip():
            i += 1
            continue

        # Paragraph (collect until blank)
        buf = [line]
        i += 1
        while i < len(lines) and lines[i].strip() and not (
            lines[i].startswith("#")
            or lines[i].startswith("```")
            or lines[i].startswith(">")
            or re.match(r"^\s*[-*+]\s+", lines[i])
            or re.match(r"^\s*\d+\.\s+", lines[i])
            or "|" in lines[i]
        ):
            buf.append(lines[i])
            i += 1
        blocks.append(Block("paragraph", text=" ".join(buf)))

    return blocks


# ---------- DOCX renderer ------------------------------------------------------

def add_runs_with_emphasis(paragraph, text):
    """Tokenize for **bold**, `code`, [text](url) links — render as runs."""
    pattern = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`|\[[^\]]+\]\([^)]+\))")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        token = m.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            run.bold = True
        elif token.startswith("`"):
            run = paragraph.add_run(token[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(10)
        elif token.startswith("["):
            link_text = re.match(r"\[([^\]]+)\]", token).group(1)
            run = paragraph.add_run(link_text)
            run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            run.font.underline = True
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def render_docx(blocks, out_path):
    doc = Document()

    # Tighten default style
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    for blk in blocks:
        if blk.kind == "heading":
            level = min(blk.level, 4)
            h = doc.add_heading(blk.text, level=level)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1F, 0x37, 0x63)

        elif blk.kind == "paragraph":
            p = doc.add_paragraph()
            add_runs_with_emphasis(p, blk.text)

        elif blk.kind == "list":
            for item in blk.items:
                # Strip leading bullet/number
                m = re.match(r"^\s*([-*+]|\d+\.)\s+(.*)", item)
                if m:
                    text = m.group(2)
                else:
                    text = item.strip()
                p = doc.add_paragraph(style="List Bullet")
                add_runs_with_emphasis(p, text)

        elif blk.kind == "code":
            # Mermaid: try to embed as PNG; fallback to source text on failure
            if blk.lang == "mermaid":
                png_path = fetch_mermaid_png(blk.content)
                if png_path is not None:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.add_run()
                    run.add_picture(str(png_path), width=Inches(6.5))
                    continue
                # else fallthrough to code rendering
            label = blk.lang.upper() if blk.lang else "CODE"
            doc.add_paragraph(f"[{label}]").runs[0].italic = True
            p = doc.add_paragraph()
            run = p.add_run(blk.content)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            # Light gray shading on the paragraph
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:fill"), "F2F2F2")
            p.paragraph_format.left_indent = Inches(0.25)
            pPr = p._p.get_or_add_pPr()
            pPr.append(shd)

        elif blk.kind == "table":
            cols = len(blk.header)
            tbl = doc.add_table(rows=1 + len(blk.rows), cols=cols)
            tbl.style = "Light Grid Accent 1"
            for j, cell in enumerate(blk.header):
                c = tbl.rows[0].cells[j]
                c.text = ""
                p = c.paragraphs[0]
                run = p.add_run(cell)
                run.bold = True
            for ri, row in enumerate(blk.rows, start=1):
                for j in range(cols):
                    val = row[j] if j < len(row) else ""
                    c = tbl.rows[ri].cells[j]
                    c.text = ""
                    p = c.paragraphs[0]
                    add_runs_with_emphasis(p, val)
                    for r in p.runs:
                        r.font.size = Pt(9)

        elif blk.kind == "quote":
            p = doc.add_paragraph()
            add_runs_with_emphasis(p, blk.content)
            p.paragraph_format.left_indent = Inches(0.5)
            for run in p.runs:
                run.italic = True
                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        elif blk.kind == "hr":
            doc.add_paragraph().add_run("─" * 60).font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    doc.save(out_path)
    print(f"DOCX written: {out_path}")


# ---------- PPTX renderer ------------------------------------------------------

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, title, summary_lines, code_excerpt=None, mermaid_png=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.word_wrap = True

    first = True
    for line in summary_lines:
        if first:
            body.text = line
            first = False
        else:
            p = body.add_paragraph()
            p.text = line
        # Compact font
        for para in body.paragraphs:
            for run in para.runs:
                run.font.size = PPt(14)

    if mermaid_png is not None:
        # Embed rendered mermaid image at the bottom of the slide
        left = PInches(0.5)
        top = PInches(4.5)
        height = PInches(2.7)
        slide.shapes.add_picture(str(mermaid_png), left, top, height=height)
    elif code_excerpt:
        # Add a code preview text box at the bottom
        left = PInches(0.5)
        top = PInches(5.2)
        width = PInches(9)
        height = PInches(2)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        tf.text = code_excerpt[:600]
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = "Consolas"
                run.font.size = PPt(10)
                run.font.color.rgb = PRGBColor(0x55, 0x55, 0x55)

    return slide


def add_diagram_slide(prs, title, mermaid_png):
    """Full-bleed diagram slide for important mermaid graphs."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    # Center the image, leaving room for title at top
    left = PInches(0.5)
    top = PInches(1.4)
    height = PInches(5.7)
    slide.shapes.add_picture(str(mermaid_png), left, top, height=height)
    return slide


def render_pptx(blocks, out_path):
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # Title slide
    add_title_slide(
        prs,
        "Deal AI — Functional Documentation",
        "Three-tier AI-driven real-estate marketing platform\nFrontend (Next.js) ▸ Backend (Express) ▸ AI Service (Flask)\nApril 2026",
    )

    # Walk blocks; group by H2 sections
    section_title = None
    section_buffer = []
    section_code = None
    section_mermaid_pngs = []

    def flush_section():
        nonlocal section_title, section_buffer, section_code, section_mermaid_pngs
        if section_title is None:
            return
        bullets = section_buffer[:]
        # If we have mermaid pngs, render dedicated diagram slides
        for png in section_mermaid_pngs:
            add_diagram_slide(prs, f"{section_title} — diagram", png)
        # Then the bullet content slides
        if not bullets and not section_code and not section_mermaid_pngs:
            bullets = ["(see full markdown for details)"]
        if bullets:
            chunks = [bullets[i : i + 7] for i in range(0, max(1, len(bullets)), 7)] or [[""]]
            for idx, chunk in enumerate(chunks):
                title = section_title if idx == 0 else f"{section_title} (cont.)"
                code = section_code if idx == 0 else None
                add_section_slide(prs, title, chunk, code)
        section_title = None
        section_buffer = []
        section_code = None
        section_mermaid_pngs = []

    for blk in blocks:
        if blk.kind == "heading":
            if blk.level == 1:
                continue  # title already covered
            if blk.level == 2:
                flush_section()
                section_title = blk.text
            else:
                # H3/H4 inserted as a bullet header
                if section_title is None:
                    section_title = blk.text
                else:
                    section_buffer.append(f"▸ {blk.text}")

        elif blk.kind == "paragraph":
            text = re.sub(r"`([^`]+)`", r"\1", blk.text)
            text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
            text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
            if len(text) > 200:
                text = text[:197] + "..."
            section_buffer.append(text)

        elif blk.kind == "list":
            for item in blk.items[:8]:
                m = re.match(r"^\s*([-*+]|\d+\.)\s+(.*)", item)
                txt = (m.group(2) if m else item.strip())
                txt = re.sub(r"`([^`]+)`", r"\1", txt)
                txt = re.sub(r"\*\*([^*]+)\*\*", r"\1", txt)
                if len(txt) > 180:
                    txt = txt[:177] + "..."
                section_buffer.append("• " + txt)

        elif blk.kind == "code":
            if blk.lang == "mermaid":
                png = fetch_mermaid_png(blk.content)
                if png is not None:
                    section_mermaid_pngs.append(png)
                    continue
                # fallback: text excerpt
            if section_code is None and blk.lang in ("mermaid", "js", "javascript", "ts",
                                                     "python", "py", "json", "bash", "sh",
                                                     "", None):
                section_code = blk.content

        elif blk.kind == "table":
            section_buffer.append(f"Table: {' / '.join(blk.header[:4])}")
            for row in blk.rows[:4]:
                snippet = " | ".join(row[:3])
                if len(snippet) > 180:
                    snippet = snippet[:177] + "..."
                section_buffer.append("  • " + snippet)
            if len(blk.rows) > 4:
                section_buffer.append(f"  • ... +{len(blk.rows) - 4} more rows")

        elif blk.kind == "quote":
            text = blk.content
            if len(text) > 180:
                text = text[:177] + "..."
            section_buffer.append("❝ " + text + " ❞")

    flush_section()

    # Closing slide
    closing = prs.slides.add_slide(prs.slide_layouts[5])
    closing.shapes.title.text = "End of deck"
    box = closing.shapes.add_textbox(PInches(1), PInches(2.5), PInches(11), PInches(3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = (
        "For full text including Mermaid diagrams, prompts, and the complete\n"
        "endpoint catalog, see DEAL_AI_FUNCTIONAL_DOCUMENTATION.md.\n\n"
        "Repos (Bits-Solutions-co):\n"
        "  • deal-ai\n"
        "  • deal-ai-server\n"
        "  • TakamolAdvancedAI"
    )
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = PPt(16)

    prs.save(out_path)
    print(f"PPTX written: {out_path} (slides: {len(prs.slides)})")


# ---------- Main ---------------------------------------------------------------

def save_with_fallback(path: Path, save_fn) -> Path:
    """Try saving to `path`; on PermissionError (file locked by Word/PowerPoint),
    fall back to a `_v2.<ext>` sibling so the output always lands somewhere."""
    try:
        save_fn(path)
        return path
    except PermissionError:
        alt = path.with_name(path.stem + "_v2" + path.suffix)
        print(f"  {path.name} is locked; writing to {alt.name} instead")
        save_fn(alt)
        return alt


def main():
    text = MD_PATH.read_text(encoding="utf-8")
    blocks = parse_md(text)
    print(f"Parsed {len(blocks)} blocks from {MD_PATH.name}")
    save_with_fallback(DOCX_PATH, lambda p: render_docx(blocks, p))
    save_with_fallback(PPTX_PATH, lambda p: render_pptx(blocks, p))


if __name__ == "__main__":
    main()
